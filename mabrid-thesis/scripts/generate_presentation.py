#!/usr/bin/env python3
"""Generate MABRID 26-slide presentation from thesis content."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# MABRID brand (from preamble.tex)
MABRID_BLUE = RGBColor(30, 58, 95)
MABRID_LIGHT = RGBColor(230, 236, 245)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(40, 40, 40)
ACCENT = RGBColor(0, 120, 180)

OUTPUT = Path(__file__).resolve().parent.parent / "MABRID_Presentation_26slides.pptx"

SLIDES: list[dict] = [
    {
        "type": "title",
        "title": "MABRID",
        "subtitle": "Place de marché numérique marocaine de confiance",
        "footer": "Découvrir. Se connecter. Grandir localement.\nRapport de projet — août 2026",
        "notes": "Accroche : MABRID n'est pas une simple app d'annonces — c'est une plateforme de confiance pour le commerce local marocain.",
    },
    {
        "type": "content",
        "title": "Le problème — Découverte fragmentée",
        "bullets": [
            "Les acheteurs peinent à trouver des entreprises locales fiables (WhatsApp, Facebook, annonces généralistes).",
            "Canaux dispersés, peu de signaux de confiance, risque perçu élevé.",
            "Les PME manquent de visibilité professionnelle au-delà des réseaux sociaux informels.",
            "Coûts d'acquisition client élevés pour les petits commerçants.",
        ],
        "notes": "Parlez 2–3 minutes avec des exemples concrets : plombier à Rabat, artisan à Marrakech, recherche sur téléphone.",
    },
    {
        "type": "content",
        "title": "Le problème — Écart de confiance",
        "bullets": [
            "Croissance illimitée sans vérification érode la confiance sur les plateformes.",
            "Littératie numérique inégale : micro-entreprises et artisans exclus.",
            "Aucun intermédiaire ne combine découverte locale + gouvernance + conformité.",
            "Opportunité : économie numérique marocaine en croissance, mobile-first.",
        ],
        "notes": "Insistez : la confiance est le frein n°1 à l'adoption du commerce numérique local.",
    },
    {
        "type": "content",
        "title": "Notre solution — MABRID",
        "bullets": [
            "Place de marché multi-acteurs : acheteurs, vendeurs locaux, prestataires de services.",
            "Découverte par catégories, carte, communautés par ville.",
            "Messagerie directe, avis structurés, modération communautaire.",
            "Application mobile MarGem + backend sécurisé (FastAPI, PostgreSQL, Azure).",
        ],
        "notes": "Solution en langage simple : une vitrine de confiance pour le commerce de proximité.",
    },
    {
        "type": "content",
        "title": "Utilisateurs cibles",
        "bullets": [
            "Acheteurs : ménages, professionnels urbains, expatriés, touristes nationaux.",
            "Vendeurs : micro-entreprises, PME, artisans, professionnels indépendants.",
            "Partenaires : chambres de commerce, associations, opérateurs télécoms.",
            "Régulateurs : CNDP, protection des consommateurs, collectivités locales.",
        ],
        "notes": "Effets réseau inter-groupes : plus de vendeurs vérifiés → plus d'acheteurs, et inversement.",
    },
    {
        "type": "content",
        "title": "Proposition de valeur",
        "bullets": [
            "Acheteurs : recherche consolidée, carte, avis authentiques, contact direct.",
            "Vendeurs : vitrine pro, gestion des annonces, analyses, badge de vérification.",
            "Plateforme : abonnements, services de confiance, partenariats entreprise.",
            "Piliers : confiance par conception · pertinence locale · inclusion · excellence opérationnelle.",
        ],
        "notes": "Différenciation vs petites annonces : vérification, modération, conformité publiée.",
    },
    {
        "type": "content",
        "title": "Démonstration — Aperçu produit",
        "bullets": [
            "Accueil : découverte par ville et catégories (alimentation, services, santé…).",
            "Vendeur : tableau de bord, annonces, messagerie, outils premium (ex. vidéo 59 s).",
            "Carte : entreprises à proximité avec profils vérifiés.",
            "Communautés : canaux par ville avec modération et signalement.",
        ],
        "notes": "Montrez 2–3 captures d'écran si disponibles. Sinon décrivez le parcours acheteur → message vendeur.",
    },
    {
        "type": "content",
        "title": "Modèle de revenus",
        "bullets": [
            "Abonnements vendeurs : niveau gratuit + niveaux Pro / Premium.",
            "Premium acheteur : découverte améliorée optionnelle (accès gratuit de base).",
            "Services de vérification et de confiance (badges, contrôles documentaires).",
            "Offres entreprise multi-sites · publicité future (politique transparente).",
            "Phase initiale : pas d'intermédiation de paiement → complexité réglementaire réduite.",
        ],
        "notes": "Flux diversifiés = résilience. CAC < CLV sur les canaux payants.",
    },
    {
        "type": "content",
        "title": "Business Model Canvas — Synthèse",
        "bullets": [
            "Segments : acheteurs urbains connectés, PME locales, entreprises multi-sites.",
            "Proposition : découverte locale de confiance avec gouvernance intégrée.",
            "Canaux : app mobile, SEO local, réseaux sociaux, terrain (chambres de commerce).",
            "Ressources clés : plateforme, modération, sécurité, marque MABRID.",
            "Partenaires : Azure, associations pro, conseil juridique, CNDP.",
        ],
        "notes": "Référez le BMC du rapport si vous avez la figure ch01-bmc.pdf.",
    },
    {
        "type": "content",
        "title": "Go-to-Market — Acquisition",
        "bullets": [
            "Acheteurs : SEO local (« meilleurs plombiers à Rabat »), ASO, influenceurs, RP confiance.",
            "Vendeurs : prospection terrain, ateliers numériques, parrainage via comptables/agences web.",
            "Entonnoir : notoriété → inscription → activation → rétention (revue hebdomadaire).",
            "Promesse de marque : « Trouvez des entreprises locales auxquelles vous pouvez faire confiance. »",
        ],
        "notes": "Phase 1 = investissement local fort avant dépenses médias nationales.",
    },
    {
        "type": "content",
        "title": "Positionnement concurrentiel",
        "bullets": [
            "Petites annonces : signaux de confiance limités, conformité souvent minimale.",
            "MABRID : vérification + avis + modération + communautés par ville.",
            "Outils vendeurs : vitrine, analyses, premium — vs annonces basiques ailleurs.",
            "Posture réglementaire : suite juridique publiée, alignement loi 09-08 / CNDP.",
        ],
        "notes": "Tableau comparatif du chapitre marketing — parlez des 4 dimensions de différenciation.",
    },
    {
        "type": "content",
        "title": "KPI commerciaux & objectifs",
        "bullets": [
            "Croissance : acheteurs et vendeurs actifs mensuels (MAU).",
            "Confiance : part de vendeurs vérifiés, délai résolution litiges.",
            "Revenus : ARPU, conversion abonnement, attrition vendeurs Pro.",
            "Court terme (0–12 mois) : lancement villes d'ancrage, premiers revenus abonnements.",
        ],
        "notes": "Reliez KPI à la gouvernance : ce qui se mesure se gère.",
    },
    {
        "type": "content",
        "title": "Structure de coûts & rentabilité",
        "bullets": [
            "Coûts cloud : calcul, stockage, sécurité (FinOps pour maîtriser la montée en charge).",
            "Personnel : ingénierie, modération, support, direction.",
            "Marketing variable par phase ; juridique et conformité semi-fixes.",
            "Seuil de rentabilité : marge brute ≥ coûts fixes — scénarios conservateur / base / optimiste.",
        ],
        "notes": "CAC < CLV. Croissance organique et parrainage améliorent les marges.",
    },
    {
        "type": "content",
        "title": "Architecture — Vue d'ensemble",
        "bullets": [
            "Client mobile Flutter (iOS/Android) → API REST + WebSockets (chat).",
            "Backend FastAPI asynchrone, PostgreSQL, authentification JWT + MFA.",
            "Stockage médias : Azure Blob (prod) ou stockage local signé (dev / homelab).",
            "Cloud principal : Microsoft Azure — Policy, RBAC, FinOps, reprise après sinistre.",
        ],
        "notes": "1 diagramme suffit. Ne lisez pas la stack — dites « conçu pour scale et sécurité ».",
    },
    {
        "type": "content",
        "title": "Architecture — Flux & déploiement",
        "bullets": [
            "Modèle hub-and-spoke : entité marocaine responsable du traitement des données.",
            "Environnements séparés : production, préproduction, gestion (micro-segmentation).",
            "Surveillance : Microsoft Sentinel, Log Analytics (SIEM/SOAR).",
            "Déploiement flexible : Azure cloud ou serveur local (pilotage / souveraineté).",
        ],
        "notes": "Mentionnez le serveur domicile comme preuve de concept, pas comme architecture finale.",
    },
    {
        "type": "content",
        "title": "Cadre juridique marocain",
        "bullets": [
            "Loi 09-08 : licéité, finalité, proportionnalité, sécurité, droits des personnes.",
            "CNDP : registre des traitements, déclarations (F211), coopération en cas d'enquête.",
            "Protection des consommateurs : annonces exactes, tarification transparente, réclamations.",
            "Commerce électronique : MABRID = plateforme intermédiaire (périmètre contractuel clair).",
        ],
        "notes": "Insistez : conformité de bonne foi + examen counsel avant lancement public à grande échelle.",
    },
    {
        "type": "content",
        "title": "Protection des données & GDPR",
        "bullets": [
            "Privacy by design : minimisation, consentement granulaire, paramètres par défaut restrictifs.",
            "Bases légales : contrat, intérêt légitime (fraude/sécurité), consentement (marketing).",
            "Droits : accès, rectification, effacement, portabilité — outils intégrés + SLA.",
            "Transferts transfrontaliers : garanties documentées (CCT, adéquation, consentement si requis).",
        ],
        "notes": "GDPR renforce la posture même pour opérations centrées Maroc si utilisateurs EEE.",
    },
    {
        "type": "content",
        "title": "Politiques de plateforme & éthique",
        "bullets": [
            "Suite publiée : CGU, confidentialité, conditions vendeurs/acheteurs, directives communautaires.",
            "Modération : contenu interdit, recours, rapports de transparence.",
            "Éthique : pas de dark patterns, annonces sponsorisées étiquetées, accessibilité WCAG.",
            "Acceptation légale tracée à l'inscription (consentement documenté, langue FR pour documents).",
        ],
        "notes": "Hiérarchie des 8 documents de politique — montrer maturité vs concurrents.",
    },
    {
        "type": "content",
        "title": "Stratégie Zero Trust",
        "bullets": [
            "Ne jamais faire confiance, toujours vérifier — au-delà du périmètre réseau.",
            "Identité d'abord : MFA obligatoire admin, RBAC (acheteur, vendeur, modérateur, admin).",
            "Moindre privilège + vérification continue (contexte, géolocalisation, risque session).",
            "Alignement NIST CSF, ISO 27001, CIS Controls — feuille de route certification.",
        ],
        "notes": "Sécurité = impératif conseil d'administration, pas une réflexion après coup.",
    },
    {
        "type": "content",
        "title": "Sécurisation applicative MarGem",
        "bullets": [
            "Auth : bcrypt, JWT courts, refresh rotatif, détection réutilisation, verrouillage compte.",
            "API : rate limiting, validation Pydantic, en-têtes sécurité, CORS strict en production.",
            "Uploads : presign, magic bytes, allowlist — validation vidéo (premium, max 59 s).",
            "WebSocket communautaire : JWT obligatoire, appartenance ville, modération.",
        ],
        "notes": "Traduction opérationnelle des principes théoriques — contrôles vérifiables.",
    },
    {
        "type": "content",
        "title": "Protection des données & incidents",
        "bullets": [
            "Classification des données, chiffrement en transit et au repos, sauvegardes testées.",
            "Gestion des menaces : évaluation des risques, modélisation, gestion des vulnérabilités.",
            "Réponse aux incidents intégrée SIEM/SOAR — playbook IR, notification CNDP si requis.",
            "Pentest indépendant et clôture des écarts avant revendication maturité opérationnelle.",
        ],
        "notes": "MTTR incidents et formation annuelle sensibilisation pour tous les employés.",
    },
    {
        "type": "content",
        "title": "Feuille de route — 3 phases",
        "bullets": [
            "Phase 1 (0–12 mois) : villes d'ancrage, vendeurs vérifiés, modérateurs locaux, KPI liquidité.",
            "Phase 2 (12–36 mois) : expansion ~20 villes, partenariats municipaux, revenus diversifiés.",
            "Phase 3 (36+ mois) : Maghreb sélectif (Tunisie, diaspora) sous faisabilité juridique.",
            "Innovations : vérification identité renforcée, intelligence économique vendeurs, API entreprise.",
        ],
        "notes": "Expansion suit le conseil réglementaire, pas le seul opportunisme.",
    },
    {
        "type": "content",
        "title": "Risques & atténuations",
        "bullets": [
            "Sanction CNDP → DPIA, DPO, registre traitements, engagement formalités.",
            "Érosion confiance → modération humaine + vérification, pas croissance volume sans qualité.",
            "Transfert données cloud → CCT, carte des flux, revue juridique, hébergement local option.",
            "Concurrence price-war → différenciation confiance, pas course au volume d'annonces.",
            "Incident sécurité → Zero Trust, SIEM, IR playbook, communication transparente.",
        ],
        "notes": "Montrez que vous avez identifié les risques — crédibilité investisseur/jury.",
    },
    {
        "type": "content",
        "title": "Vision",
        "bullets": [
            "Devenir la place de marché numérique la plus fiable du Maroc.",
            "Mission : découverte transparente, signaux de confiance vérifiables, gouvernance responsable.",
            "Impact : numérisation des PME, protection des consommateurs, emploi local.",
            "Succès = confiance par gouvernance, sécurité et intégrité juridique — pas seulement effets de réseau.",
        ],
        "notes": "Vision inspirante mais ancrée dans exécution disciplinée.",
    },
    {
        "type": "content",
        "title": "Demande & prochaines étapes",
        "bullets": [
            "Examen juridique formel + enregistrement CNDP avant lancement public à grande échelle.",
            "Évaluation sécurité indépendante + plan clôture écarts ISO 27001.",
            "Pilotage villes d'ancrage : cohorte vendeurs vérifiés + métriques transparence.",
            "Partenariats : chambres de commerce, financement seed, co-marketing entreprise.",
            "Ressources : équipe modération, budget marketing local, infrastructure cloud maîtrisée.",
        ],
        "notes": "Adaptez l'ask à l'audience : jury = méthodologie ; investisseur = montant + jalons.",
    },
    {
        "type": "closing",
        "title": "Merci",
        "subtitle": "Questions & échanges",
        "footer": "[Nom de l'auteur]\n[Email] · [LinkedIn]\nMABRID — Maroc, août 2026",
        "notes": "Terminez par la promesse : « Trouvez des entreprises locales auxquelles vous pouvez faire confiance. »",
    },
]


def _set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_header_bar(slide, prs: Presentation) -> None:
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0),
        Inches(0),
        prs.slide_width,
        Inches(0.12),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MABRID_BLUE
    shape.line.fill.background()


def _add_footer(slide, prs: Presentation, text: str = "MABRID — Place de marché de confiance") -> None:
    box = slide.shapes.add_textbox(
        Inches(0.5), prs.slide_height - Inches(0.45), prs.slide_width - Inches(1), Inches(0.35)
    )
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(120, 120, 120)
    p.alignment = PP_ALIGN.CENTER


def _style_title(text_frame, size: int = 32, color: RGBColor = MABRID_BLUE) -> None:
    for i, para in enumerate(text_frame.paragraphs):
        para.font.bold = True
        para.font.size = Pt(size if i == 0 else size - 4)
        para.font.color.rgb = color
        para.font.name = "Calibri"


def _add_bullets(text_frame, bullets: list[str]) -> None:
    text_frame.clear()
    for i, bullet in enumerate(bullets):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.font.name = "Calibri"
        p.space_after = Pt(10)


def _add_notes(slide, notes: str) -> None:
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    for spec in SLIDES:
        slide = prs.slides.add_slide(blank)
        stype = spec["type"]

        if stype == "title":
            _set_slide_bg(slide, MABRID_BLUE)
            _add_header_bar(slide, prs)

            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.5)
            )
            tf = title_box.text_frame
            tf.text = spec["title"]
            tf.paragraphs[0].font.size = Pt(54)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = WHITE
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            sub_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(3.5), Inches(11.3), Inches(1.2)
            )
            stf = sub_box.text_frame
            stf.text = spec.get("subtitle", "")
            stf.paragraphs[0].font.size = Pt(24)
            stf.paragraphs[0].font.color.rgb = MABRID_LIGHT
            stf.paragraphs[0].alignment = PP_ALIGN.CENTER

            foot_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.0)
            )
            ftf = foot_box.text_frame
            ftf.text = spec.get("footer", "")
            for para in ftf.paragraphs:
                para.font.size = Pt(14)
                para.font.color.rgb = RGBColor(200, 210, 225)
                para.alignment = PP_ALIGN.CENTER

        elif stype == "section":
            _set_slide_bg(slide, MABRID_LIGHT)
            _add_header_bar(slide, prs)

            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.5)
            )
            tf = title_box.text_frame
            tf.text = spec["title"]
            _style_title(tf, size=40)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            if spec.get("subtitle"):
                sub = tf.add_paragraph()
                sub.text = spec["subtitle"]
                sub.font.size = Pt(22)
                sub.font.color.rgb = ACCENT
                sub.alignment = PP_ALIGN.CENTER

        elif stype == "closing":
            _set_slide_bg(slide, MABRID_BLUE)
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.2)
            )
            tf = title_box.text_frame
            tf.text = spec["title"]
            tf.paragraphs[0].font.size = Pt(48)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = WHITE
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            sub_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(3.3), Inches(11.3), Inches(0.8)
            )
            stf = sub_box.text_frame
            stf.text = spec.get("subtitle", "")
            stf.paragraphs[0].font.size = Pt(22)
            stf.paragraphs[0].font.color.rgb = MABRID_LIGHT
            stf.paragraphs[0].alignment = PP_ALIGN.CENTER

            foot_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.5)
            )
            ftf = foot_box.text_frame
            ftf.text = spec.get("footer", "")
            for para in ftf.paragraphs:
                para.font.size = Pt(16)
                para.font.color.rgb = RGBColor(200, 210, 225)
                para.alignment = PP_ALIGN.CENTER

        else:  # content
            _set_slide_bg(slide, WHITE)
            _add_header_bar(slide, prs)
            _add_footer(slide, prs)

            title_box = slide.shapes.add_textbox(
                Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9)
            )
            tf = title_box.text_frame
            tf.text = spec["title"]
            _style_title(tf, size=30)
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

            body_box = slide.shapes.add_textbox(
                Inches(0.7), Inches(1.35), Inches(12.0), Inches(5.5)
            )
            btf = body_box.text_frame
            btf.word_wrap = True
            _add_bullets(btf, spec.get("bullets", []))

        _add_notes(slide, spec.get("notes", ""))

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build_presentation()
    print(f"Generated: {path}")
    print(f"Slides: {len(SLIDES)}")
